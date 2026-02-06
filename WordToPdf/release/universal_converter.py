"""
通用文件格式转换器
支持 PDF, Word, HTML, PPT, 图片(PNG/JPG/BMP), CSV 之间的互转
"""

import os
import sys
import subprocess
import tempfile
from pathlib import Path
from typing import Optional, List, Dict, Any
import logging

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 导入转换库
try:
    import pandas as pd
    from PIL import Image
    from docx import Document
    import python_pptx as pptx
    from pptx import Presentation
    from bs4 import BeautifulSoup
    import img2pdf
    from fpdf import FPDF
    import pdf2image
    import pdfplumber
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.utils import ImageReader
except ImportError as e:
    logger.error(f"缺少必要的依赖库: {e}")
    logger.error("请安装所有依赖: pip install pandas openpyxl python-docx pillow python-pptx beautifulsoup4 html5lib img2pdf pdf2image pdfplumber reportlab fpdf2")
    sys.exit(1)

# PDF相关库
try:
    from pdf2docx import Converter
    from PyPDF2 import PdfReader, PdfWriter
except ImportError as e:
    logger.warning(f"PDF相关库未完全安装: {e}")


class UniversalConverter:
    """通用文件格式转换器"""
    
    # 支持的文件格式
    SUPPORTED_FORMATS = {
        'pdf': ['.pdf'],
        'word': ['.docx', '.doc'],
        'html': ['.html', '.htm'],
        'ppt': ['.pptx', '.ppt'],
        'image': ['.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff'],
        'csv': ['.csv', '.xlsx', '.xls']
    }
    
    def __init__(self):
        """初始化转换器"""
        self.temp_dir = tempfile.gettempdir()
        
    def get_file_type(self, file_path: str) -> Optional[str]:
        """根据文件扩展名获取文件类型"""
        ext = Path(file_path).suffix.lower()
        for file_type, extensions in self.SUPPORTED_FORMATS.items():
            if ext in extensions:
                return file_type
        return None
    
    def validate_file(self, file_path: str) -> bool:
        """验证文件是否存在且有效"""
        if not os.path.exists(file_path):
            logger.error(f"文件不存在: {file_path}")
            return False
        
        if not os.path.isfile(file_path):
            logger.error(f"路径不是有效文件: {file_path}")
            return False
        
        file_type = self.get_file_type(file_path)
        if not file_type:
            logger.error(f"不支持的文件格式: {file_path}")
            return False
        
        return True
    
    def convert(self, input_path: str, output_path: str) -> bool:
        """
        主转换方法
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            
        Returns:
            bool: 转换是否成功
        """
        try:
            # 验证输入文件
            if not self.validate_file(input_path):
                return False
            
            # 获取输入和输出类型
            input_type = self.get_file_type(input_path)
            output_type = self.get_file_type(output_path)
            
            if not input_type or not output_type:
                logger.error("无法识别文件格式")
                return False
            
            # 创建输出目录
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            logger.info(f"开始转换: {input_type} -> {output_type}")
            
            # 根据转换类型调用相应方法
            conversion_method = f"_convert_{input_type}_to_{output_type}"
            method = getattr(self, conversion_method, None)
            
            if method:
                return method(input_path, output_path)
            else:
                logger.error(f"不支持的转换: {input_type} -> {output_type}")
                return False
                
        except Exception as e:
            logger.error(f"转换失败: {e}")
            return False
    
    def convert_batch(self, input_files: List[str], output_dir: str, output_format: str) -> Dict[str, bool]:
        """
        批量转换文件
        
        Args:
            input_files: 输入文件列表
            output_dir: 输出目录
            output_format: 输出格式 (如 'pdf', 'word', 等)
            
        Returns:
            Dict[str, bool]: 转换结果字典 {文件名: 是否成功}
        """
        results = {}
        
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        for input_file in input_files:
            if not self.validate_file(input_file):
                results[input_file] = False
                continue
            
            # 生成输出文件名
            input_name = Path(input_file).stem
            output_ext = self.SUPPORTED_FORMATS.get(output_format, [''])[0]
            output_path = os.path.join(output_dir, f"{input_name}{output_ext}")
            
            # 转换文件
            success = self.convert(input_file, output_path)
            results[input_file] = success
        
        return results
    
    # PDF转换方法
    def _convert_pdf_to_word(self, input_path: str, output_path: str) -> bool:
        """PDF转Word"""
        try:
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            logger.info("PDF转Word成功")
            return True
        except Exception as e:
            logger.error(f"PDF转Word失败: {e}")
            return False
    
    def _convert_pdf_to_image(self, input_path: str, output_path: str) -> bool:
        """PDF转图片"""
        try:
            # 确保输出路径支持图片格式
            output_ext = Path(output_path).suffix.lower()
            if output_ext not in ['.png', '.jpg', '.jpeg', '.bmp']:
                output_path = os.path.splitext(output_path)[0] + '.png'
            
            images = pdf2image.convert_from_path(input_path)
            if len(images) == 1:
                images[0].save(output_path)
            else:
                # 多页PDF，每页保存为单独文件
                base_name = os.path.splitext(output_path)[0]
                for i, image in enumerate(images):
                    image.save(f"{base_name}_{i+1}{os.path.splitext(output_path)[1]}")
            
            logger.info("PDF转图片成功")
            return True
        except Exception as e:
            logger.error(f"PDF转图片失败: {e}")
            return False
    
    def _convert_pdf_to_html(self, input_path: str, output_path: str) -> bool:
        """PDF转HTML"""
        try:
            # 使用pdfplumber提取文本和表格
            html_content = "<html><head><meta charset='UTF-8'><title>PDF转HTML</title></head><body>"
            
            with pdfplumber.open(input_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    html_content += f"<h2>第{page_num+1}页</h2>"
                    text = page.extract_text()
                    if text:
                        html_content += f"<pre>{text}</pre>"
                    
                    # 提取表格
                    tables = page.extract_tables()
                    for table in tables:
                        html_content += "<table border='1'>"
                        for row in table:
                            html_content += "<tr>"
                            for cell in row:
                                html_content += f"<td>{cell if cell else ''}</td>"
                            html_content += "</tr>"
                        html_content += "</table>"
            
            html_content += "</body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            logger.info("PDF转HTML成功")
            return True
        except Exception as e:
            logger.error(f"PDF转HTML失败: {e}")
            return False
    
    def _convert_pdf_to_csv(self, input_path: str, output_path: str) -> bool:
        """PDF转CSV"""
        try:
            # 尝试提取表格数据
            all_data = []
            
            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        # 清理表格数据
                        clean_table = []
                        for row in table:
                            clean_row = [str(cell).strip() if cell else '' for cell in row]
                            clean_table.append(clean_row)
                        all_data.extend(clean_table)
            
            if all_data:
                df = pd.DataFrame(all_data)
                df.to_csv(output_path, index=False, encoding='utf-8-sig')
                logger.info("PDF转CSV成功")
                return True
            else:
                # 如果没有表格，提取文本保存为CSV
                with pdfplumber.open(input_path) as pdf:
                    text_data = []
                    for page_num, page in enumerate(pdf.pages):
                        text = page.extract_text()
                        if text:
                            lines = text.split('\n')
                            for line in lines:
                                text_data.append([page_num+1, line.strip()])
                    
                    if text_data:
                        df = pd.DataFrame(text_data, columns=['Page', 'Text'])
                        df.to_csv(output_path, index=False, encoding='utf-8-sig')
                        logger.info("PDF转CSV成功（文本模式）")
                        return True
            
            logger.error("PDF中未找到可转换为CSV的数据")
            return False
        except Exception as e:
            logger.error(f"PDF转CSV失败: {e}")
            return False
    
    # Word转换方法
    def _convert_word_to_pdf(self, input_path: str, output_path: str) -> bool:
        """Word转PDF"""
        try:
            # 尝试使用docx2pdf（依赖Word）
            from docx2pdf import convert as docx_to_pdf
            docx_to_pdf(input_path, output_path)
            logger.info("Word转PDF成功")
            return True
        except:
            # 使用备用方案
            try:
                doc = Document(input_path)
                pdf = FPDF()
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                
                for para in doc.paragraphs:
                    if para.text.strip():
                        pdf.cell(0, 10, para.text, ln=1)
                
                pdf.output(output_path)
                logger.info("Word转PDF成功（备用方案）")
                return True
            except Exception as e:
                logger.error(f"Word转PDF失败: {e}")
                return False
    
    def _convert_word_to_html(self, input_path: str, output_path: str) -> bool:
        """Word转HTML"""
        try:
            doc = Document(input_path)
            html_content = """<html>
<head>
    <meta charset="UTF-8">
    <title>Word转HTML</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 40px; }
        p { margin-bottom: 16px; line-height: 1.5; }
        h1, h2, h3 { margin-top: 24px; margin-bottom: 12px; }
        table { border-collapse: collapse; width: 100%; margin-bottom: 20px; }
        td, th { border: 1px solid #ddd; padding: 8px; text-align: left; }
        th { background-color: #f2f2f2; }
    </style>
</head>
<body>
"""
            
            for element in doc.element.body:
                if element.tag.endswith('p'):
                    # 段落
                    para = element
                    text = ''.join(node.text for node in para.iter() if hasattr(node, 'text'))
                    html_content += f"<p>{text}</p>"
                elif element.tag.endswith('tbl'):
                    # 表格
                    html_content += "<table>"
                    for row in element.iter():
                        if row.tag.endswith('tr'):
                            html_content += "<tr>"
                            for cell in row.iter():
                                if cell.tag.endswith('tc') or cell.tag.endswith('p'):
                                    text = ''.join(node.text for node in cell.iter() if hasattr(node, 'text'))
                                    html_content += f"<td>{text}</td>"
                            html_content += "</tr>"
                    html_content += "</table>"
            
            html_content += "</body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            logger.info("Word转HTML成功")
            return True
        except Exception as e:
            logger.error(f"Word转HTML失败: {e}")
            return False
    
    def _convert_word_to_csv(self, input_path: str, output_path: str) -> bool:
        """Word转CSV"""
        try:
            doc = Document(input_path)
            data = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    data.append([para.text.strip()])
            
            if data:
                df = pd.DataFrame(data, columns=['Text'])
                df.to_csv(output_path, index=False, encoding='utf-8-sig')
                logger.info("Word转CSV成功")
                return True
            else:
                logger.error("Word文档中没有找到文本内容")
                return False
        except Exception as e:
            logger.error(f"Word转CSV失败: {e}")
            return False
    
    # HTML转换方法
    def _convert_html_to_pdf(self, input_path: str, output_path: str) -> bool:
        """HTML转PDF"""
        try:
            # 读取HTML内容
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 使用BeautifulSoup解析HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            text = soup.get_text()
            
            # 创建PDF
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            # 分行写入文本
            lines = text.split('\n')
            for line in lines:
                if line.strip():
                    pdf.cell(0, 10, line.strip(), ln=1)
            
            pdf.output(output_path)
            logger.info("HTML转PDF成功")
            return True
        except Exception as e:
            logger.error(f"HTML转PDF失败: {e}")
            return False
    
    def _convert_html_to_word(self, input_path: str, output_path: str) -> bool:
        """HTML转Word"""
        try:
            # 读取HTML内容
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # 使用BeautifulSoup解析HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 创建Word文档
            doc = Document()
            
            # 提取文本并添加到Word文档
            text = soup.get_text()
            lines = text.split('\n')
            
            for line in lines:
                if line.strip():
                    doc.add_paragraph(line.strip())
            
            doc.save(output_path)
            logger.info("HTML转Word成功")
            return True
        except Exception as e:
            logger.error(f"HTML转Word失败: {e}")
            return False
    
    # PPT转换方法
    def _convert_ppt_to_pdf(self, input_path: str, output_path: str) -> bool:
        """PPT转PDF"""
        try:
            # 尝试使用LibreOffice (需要系统安装)
            command = [
                'soffice', '--headless', '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_path),
                input_path
            ]
            
            result = subprocess.run(command, capture_output=True, text=True)
            if result.returncode == 0:
                logger.info("PPT转PDF成功")
                return True
            else:
                # 备用方案：将每页转换为图片再转PDF
                return self._convert_ppt_to_image_to_pdf(input_path, output_path)
        except Exception as e:
            logger.error(f"PPT转PDF失败: {e}")
            return False
    
    def _convert_ppt_to_image(self, input_path: str, output_path: str) -> bool:
        """PPT转图片"""
        try:
            prs = Presentation(input_path)
            
            # 确保输出路径支持图片格式
            output_ext = Path(output_path).suffix.lower()
            if output_ext not in ['.png', '.jpg', '.jpeg', '.bmp']:
                output_path = os.path.splitext(output_path)[0] + '.png'
            
            base_name = os.path.splitext(output_path)[0]
            ext = os.path.splitext(output_path)[1]
            
            for i, slide in enumerate(prs.slides):
                # 创建临时图片
                img_path = f"{base_name}_{i+1}{ext}"
                
                # 这里需要使用额外的库来渲染PPT为图片
                # 由于PPTX库不支持直接渲染，使用占位符方法
                logger.info(f"处理第{i+1}页...")
            
            logger.info("PPT转图片功能需要额外依赖")
            return False
        except Exception as e:
            logger.error(f"PPT转图片失败: {e}")
            return False
    
    def _convert_ppt_to_image_to_pdf(self, input_path: str, output_path: str) -> bool:
        """PPT转PDF的备用方案"""
        try:
            # 这是一个简化的实现，实际需要更复杂的逻辑
            prs = Presentation(input_path)
            pdf = FPDF()
            
            for slide in prs.slides:
                pdf.add_page()
                pdf.set_font("Arial", size=12)
                
                # 提取幻灯片文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pdf.cell(0, 10, shape.text, ln=1)
            
            pdf.output(output_path)
            logger.info("PPT转PDF成功（文本模式）")
            return True
        except Exception as e:
            logger.error(f"PPT转PDF失败: {e}")
            return False
    
    # 图片转换方法
    def _convert_image_to_pdf(self, input_path: str, output_path: str) -> bool:
        """图片转PDF"""
        try:
            # 使用img2pdf库
            with open(output_path, "wb") as f:
                f.write(img2pdf.convert(input_path))
            
            logger.info("图片转PDF成功")
            return True
        except Exception as e:
            # 备用方案
            try:
                image = Image.open(input_path)
                pdf = FPDF()
                pdf.add_page()
                
                # 获取图片尺寸
                width, height = image.size
                pdf_w, pdf_h = pdf.w, pdf.h
                
                # 计算缩放比例
                scale = min(pdf_w/width, pdf_h/height) * 0.8
                new_width = width * scale
                new_height = height * scale
                
                x = (pdf_w - new_width) / 2
                y = (pdf_h - new_height) / 2
                
                pdf.image(input_path, x, y, new_width, new_height)
                pdf.output(output_path)
                
                logger.info("图片转PDF成功（备用方案）")
                return True
            except Exception as e2:
                logger.error(f"图片转PDF失败: {e2}")
                return False
    
    def _convert_image_to_html(self, input_path: str, output_path: str) -> bool:
        """图片转HTML"""
        try:
            html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>图片转HTML</title>
    <style>
        body {{ margin: 0; padding: 20px; text-align: center; }}
        img {{ max-width: 100%; height: auto; }}
    </style>
</head>
<body>
    <h1>图片预览</h1>
    <img src="{os.path.basename(input_path)}" alt="转换的图片">
</body>
</html>"""
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            # 复制图片到输出目录
            output_dir = os.path.dirname(output_path)
            if output_dir:
                import shutil
                shutil.copy2(input_path, output_dir)
            
            logger.info("图片转HTML成功")
            return True
        except Exception as e:
            logger.error(f"图片转HTML失败: {e}")
            return False
    
    # CSV转换方法
    def _convert_csv_to_pdf(self, input_path: str, output_path: str) -> bool:
        """CSV转PDF"""
        try:
            df = pd.read_csv(input_path, encoding='utf-8-sig')
            
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            
            # 写入表头
            for col in df.columns:
                pdf.cell(40, 10, str(col), border=1)
            pdf.ln()
            
            # 写入数据行
            for index, row in df.iterrows():
                for col in df.columns:
                    pdf.cell(40, 10, str(row[col]), border=1)
                pdf.ln()
            
            pdf.output(output_path)
            logger.info("CSV转PDF成功")
            return True
        except Exception as e:
            logger.error(f"CSV转PDF失败: {e}")
            return False
    
    def _convert_csv_to_html(self, input_path: str, output_path: str) -> bool:
        """CSV转HTML"""
        try:
            df = pd.read_csv(input_path, encoding='utf-8-sig')
            
            html_content = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>CSV转HTML</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 20px; }
        table { border-collapse: collapse; width: 100%; }
        th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
        th { background-color: #f2f2f2; }
        tr:nth-child(even) { background-color: #f9f9f9; }
    </style>
</head>
<body>
    <h1>CSV数据表格</h1>
""" + df.to_html(index=False, escape=False) + """
</body>
</html>"""
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            logger.info("CSV转HTML成功")
            return True
        except Exception as e:
            logger.error(f"CSV转HTML失败: {e}")
            return False
    
    def _convert_csv_to_excel(self, input_path: str, output_path: str) -> bool:
        """CSV转Excel"""
        try:
            df = pd.read_csv(input_path, encoding='utf-8-sig')
            df.to_excel(output_path, index=False)
            logger.info("CSV转Excel成功")
            return True
        except Exception as e:
            logger.error(f"CSV转Excel失败: {e}")
            return False
    
    def get_supported_conversions(self) -> Dict[str, List[str]]:
        """获取支持的转换格式"""
        conversions = {}
        for from_type in self.SUPPORTED_FORMATS.keys():
            conversions[from_type] = []
            for to_type in self.SUPPORTED_FORMATS.keys():
                if from_type != to_type:
                    method_name = f"_convert_{from_type}_to_{to_type}"
                    if hasattr(self, method_name):
                        conversions[from_type].append(to_type)
        return conversions